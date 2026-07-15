"""Extraction: native PDF / DOCX text with bounding boxes, OCR routing for scans.

Every extracted span keeps (page_number, x0, y0, x1, y1, text) -- these
coordinates power the citation-highlight feature, so they are preserved all
the way into chunk metadata.

Coordinate space: PDF points as reported by PyMuPDF (origin top-left). The
page-image endpoint reports page width/height in the same space so the
frontend can scale highlight rectangles to the rendered image.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, List, Optional

import fitz  # PyMuPDF

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".log"}
ODF_EXTS = {".odt", ".ods", ".odp"}
# Pre-2007 binary Office formats: tracked so the user gets a clear
# "re-save as modern format" message instead of silent absence.
LEGACY_EXTS = {".doc", ".xls", ".ppt"}
SUPPORTED_EXTS = (
    {".pdf", ".docx", ".xlsx", ".pptx", ".html", ".htm", ".rtf", ".eml"}
    | TEXT_EXTS | ODF_EXTS | LEGACY_EXTS | IMAGE_EXTS
)


class UnsupportedFormatError(Exception):
    """Raised with a user-facing explanation of why a file can't be indexed."""

# A page whose native text layer has fewer chars than this is treated as scanned.
MIN_NATIVE_CHARS_PER_PAGE = 24


@dataclass
class Span:
    page_number: int  # 0-based
    x0: float
    y0: float
    x1: float
    y1: float
    text: str
    is_heading: bool = False


@dataclass
class ExtractionResult:
    spans: List[Span] = field(default_factory=list)
    page_count: int = 0
    used_ocr: bool = False
    ocr_unavailable_pages: int = 0  # scanned pages we could not OCR


ProgressCb = Callable[[str], None]


def extract(path: str, progress: Optional[ProgressCb] = None) -> ExtractionResult:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path, progress)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext in (".html", ".htm"):
        return _extract_html(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    if ext == ".eml":
        return _extract_eml(path)
    if ext in ODF_EXTS:
        return _extract_odf(path)
    if ext in TEXT_EXTS:
        return _extract_text(path)
    if ext in IMAGE_EXTS:
        return _extract_image(path, progress)
    if ext in LEGACY_EXTS:
        modern = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}[ext]
        raise UnsupportedFormatError(
            f"Old Office format ({ext}). Open it in Word/Excel/PowerPoint and "
            f"save as {modern} — it will then be indexed automatically."
        )
    raise UnsupportedFormatError(f"Unsupported file type: {ext}")


# --- PDF ---------------------------------------------------------------------


def _extract_pdf(path: str, progress: Optional[ProgressCb]) -> ExtractionResult:
    result = ExtractionResult()
    doc = fitz.open(path)
    result.page_count = doc.page_count
    median_font_sizes: List[float] = []

    for pno in range(doc.page_count):
        page = doc[pno]
        d = page.get_text("dict")
        page_spans, font_sizes = _spans_from_textdict(d, pno)
        native_chars = sum(len(s.text) for s in page_spans)

        if native_chars >= MIN_NATIVE_CHARS_PER_PAGE:
            result.spans.extend(page_spans)
            median_font_sizes.extend(font_sizes)
        else:
            # Scanned page -> OCR route.
            if progress:
                progress(f"OCR: page {pno + 1} of {doc.page_count}")
            from .ocr import ocr_pdf_page

            ocr_spans = ocr_pdf_page(page, pno)
            if ocr_spans is None:
                result.ocr_unavailable_pages += 1
            else:
                result.used_ocr = True
                result.spans.extend(ocr_spans)

    _mark_headings(result.spans, median_font_sizes)
    doc.close()
    return result


def _spans_from_textdict(d: dict, pno: int):
    spans: List[Span] = []
    font_sizes: List[float] = []
    for block in d.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            text = "".join(s.get("text", "") for s in line.get("spans", []))
            if not text.strip():
                continue
            x0, y0, x1, y1 = line["bbox"]
            sizes = [s.get("size", 0) for s in line.get("spans", [])]
            size = max(sizes) if sizes else 0.0
            bold = any("bold" in s.get("font", "").lower() for s in line.get("spans", []))
            sp = Span(pno, x0, y0, x1, y1, text)
            sp.is_heading = bold and len(text.strip()) < 90
            spans.append(sp)
            font_sizes.append(size)
    return spans, font_sizes


def _mark_headings(spans: List[Span], font_sizes: List[float]) -> None:
    """Also flag lines noticeably larger than the median font as headings."""
    if not font_sizes:
        return
    med = sorted(font_sizes)[len(font_sizes) // 2]
    # is_heading from bold detection stays; size-based pass handled in
    # _spans_from_textdict would need per-span size, so we keep bold+short only
    # plus a cheap ALL-CAPS heuristic here.
    for sp in spans:
        t = sp.text.strip()
        if not sp.is_heading and 3 < len(t) < 80 and t.isupper():
            sp.is_heading = True
    _ = med  # median kept for future size-based heading detection


# --- DOCX --------------------------------------------------------------------


def _extract_docx(path: str) -> ExtractionResult:
    """DOCX has no fixed page geometry; we emit spans with a synthetic layout
    (one 'page', stacked lines) so citations still resolve to a text location."""
    import docx

    result = ExtractionResult()
    document = docx.Document(path)
    y = 0.0
    line_h = 14.0
    page_h = 792.0  # letter @ 72dpi
    pno = 0
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if y + line_h > page_h:
            pno += 1
            y = 0.0
        sp = Span(pno, 36.0, y, 576.0, y + line_h, text)
        style = (para.style.name or "").lower() if para.style else ""
        sp.is_heading = style.startswith("heading") or style == "title"
        result.spans.append(sp)
        y += line_h * (1.6 if sp.is_heading else 1.15)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if not cells:
                continue
            if y + line_h > page_h:
                pno += 1
                y = 0.0
            result.spans.append(Span(pno, 36.0, y, 576.0, y + line_h, " | ".join(cells)))
            y += line_h * 1.15
    result.page_count = pno + 1
    return result


# --- text-like formats (txt/md/csv/html/xlsx/pptx) ---------------------------
# These have no fixed page geometry, so like DOCX they get a synthetic layout;
# citations for them open the highlighted-text view rather than a page image.

MAX_TEXT_LINES = 20000  # guard against enormous CSV/log files


def _spans_from_lines(lines: List[tuple]) -> ExtractionResult:
    """lines: [(text, is_heading)] -> spans with a synthetic stacked layout."""
    result = ExtractionResult()
    y, line_h, page_h, pno = 0.0, 14.0, 792.0, 0
    for text, is_heading in lines[:MAX_TEXT_LINES]:
        text = text.strip()
        if not text:
            continue
        if y + line_h > page_h:
            pno += 1
            y = 0.0
        sp = Span(pno, 36.0, y, 576.0, y + line_h, text)
        sp.is_heading = is_heading
        result.spans.append(sp)
        y += line_h * (1.6 if is_heading else 1.15)
    result.page_count = pno + 1
    return result


def _extract_text(path: str) -> ExtractionResult:
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    is_md = Path(path).suffix.lower() == ".md"
    lines = []
    for line in raw.splitlines():
        heading = is_md and line.lstrip().startswith("#")
        lines.append((line.lstrip("# ") if heading else line, heading))
    return _spans_from_lines(lines)


def _extract_html(path: str) -> ExtractionResult:
    from html.parser import HTMLParser

    class _Text(HTMLParser):
        def __init__(self):
            super().__init__()
            self.lines: List[tuple] = []
            self._skip = 0
            self._heading = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip += 1
            if tag in ("h1", "h2", "h3", "h4"):
                self._heading += 1

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self._skip:
                self._skip -= 1
            if tag in ("h1", "h2", "h3", "h4") and self._heading:
                self._heading -= 1

        def handle_data(self, data):
            if not self._skip and data.strip():
                self.lines.append((data.strip(), self._heading > 0))

    parser = _Text()
    parser.feed(Path(path).read_text(encoding="utf-8", errors="replace"))
    return _spans_from_lines(parser.lines)


def _extract_xlsx(path: str) -> ExtractionResult:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    lines: List[tuple] = []
    for ws in wb.worksheets:
        lines.append((f"Sheet: {ws.title}", True))
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append((" | ".join(cells), False))
            if len(lines) > MAX_TEXT_LINES:
                break
    wb.close()
    return _spans_from_lines(lines)


def _extract_pptx(path: str) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(path)
    result = ExtractionResult()
    line_h = 14.0
    for slide_no, slide in enumerate(prs.slides):
        y = 0.0
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                sp = Span(slide_no, 36.0, y, 576.0, y + line_h, text)
                # Title placeholders read as headings.
                sp.is_heading = bool(
                    shape == slide.shapes.title if slide.shapes.title else False
                )
                result.spans.append(sp)
                y += line_h * 1.2
    result.page_count = max(1, len(prs.slides))
    return result


def _extract_rtf(path: str) -> ExtractionResult:
    from striprtf.striprtf import rtf_to_text

    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw, errors="ignore")
    return _spans_from_lines([(line, False) for line in text.splitlines()])


def _extract_eml(path: str) -> ExtractionResult:
    import email
    from email import policy

    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=policy.default)
    lines: List[tuple] = [
        (f"Subject: {msg.get('Subject', '(no subject)')}", True),
        (f"From: {msg.get('From', '')}  To: {msg.get('To', '')}  Date: {msg.get('Date', '')}", False),
    ]
    body = msg.get_body(preferencelist=("plain", "html"))
    if body is not None:
        content = body.get_content()
        if body.get_content_type() == "text/html":
            import re as _re

            content = _re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", content, flags=_re.S | _re.I)
            content = _re.sub(r"<[^>]+>", " ", content)
        for line in content.splitlines():
            lines.append((line, False))
    return _spans_from_lines(lines)


def _extract_odf(path: str) -> ExtractionResult:
    """OpenDocument (odt/ods/odp): pull text from content.xml with stdlib."""
    import re as _re
    import zipfile

    with zipfile.ZipFile(path) as zf:
        xml = zf.read("content.xml").decode("utf-8", errors="replace")
    lines: List[tuple] = []
    # Paragraphs and headings carry the visible text in all three ODF types.
    for m in _re.finditer(r"<text:(h|p)\b[^>]*>(.*?)</text:\1>", xml, flags=_re.S):
        kind, inner = m.group(1), m.group(2)
        text = _re.sub(r"<[^>]+>", " ", inner)
        text = _re.sub(r"\s+", " ", text).strip()
        if text:
            lines.append((text, kind == "h"))
    return _spans_from_lines(lines)


# --- standalone images -------------------------------------------------------


def _extract_image(path: str, progress: Optional[ProgressCb]) -> ExtractionResult:
    from .ocr import ocr_image_file

    result = ExtractionResult(page_count=1)
    if progress:
        progress("OCR: image")
    spans = ocr_image_file(path)
    if spans is None:
        result.ocr_unavailable_pages = 1
    else:
        result.used_ocr = True
        result.spans = spans
    return result
